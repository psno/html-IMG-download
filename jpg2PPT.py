import os
import re
from pptx import Presentation
# 注意：Inches 在此脚本中未被使用，可以移除
from PIL import Image
import tempfile

# --- 配置 ---
# 定义标准幻灯片的宽高比 (例如 16:9)
STANDARD_ASPECT_RATIO = 16 / 9
# 定义容差，用于判断实际分割高度与理想高度的接近程度 (例如 10%)
TOLERANCE_PERCENTAGE = 0.15
# --- 配置结束 ---


def detect_and_split_vertical_images(image_path, aspect_ratio=STANDARD_ASPECT_RATIO, tolerance=TOLERANCE_PERCENTAGE):
    """
    检测图片是否为纵向拼接的多张图片，如果是则分割。
    通过图片宽度和预设宽高比动态估算单页高度。

    Args:
        image_path: 图片路径
        aspect_ratio: 标准单张幻灯片的宽高比 (width/height)，默认为 16/9。
        tolerance: 容差百分比，用于判断分割后高度是否合理，默认为 0.15 (15%)。

    Returns:
        分割后的图片路径列表，如果不需要分割则返回包含原图片路径的列表。
    """
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            print(f"  图片尺寸: {width} x {height}")

            # 如果图片宽度大于高度，很可能不是纵向拼接图，直接返回
            if width >= height:
                 print(f"  图片宽度 >= 高度，跳过分割检查。")
                 return [image_path]

            # 基于图片宽度和标准宽高比计算理想的单页高度
            ideal_page_height = width / aspect_ratio
            print(f"  理想单页高度 (基于宽高比 {aspect_ratio:.2f}): {ideal_page_height:.2f}")

            # 计算可能的分割数量
            split_count = round(height / ideal_page_height)
            print(f"  初步估算分割数量: {split_count}")

            # 如果估算的分割数量小于等于1，则不进行分割
            if split_count <= 1:
                print(f"  估算分割数量 <= 1，无需分割。")
                return [image_path]

            # 计算实际分割后的单页高度
            actual_page_height = height / split_count
            print(f"  实际单页高度: {actual_page_height:.2f}")

            # 检查实际高度与理想高度的差异是否在容差范围内
            height_diff_ratio = abs(actual_page_height - ideal_page_height) / ideal_page_height
            print(f"  高度差异比率: {height_diff_ratio:.2%}")

            # 如果差异在容差内，则认为是拼接图片并进行分割
            if height_diff_ratio <= tolerance:
                print(f"  高度差异比率 <= {tolerance:.0%}，确认为拼接图片，将分割为 {split_count} 张")
                
                # 创建临时目录存储分割后的图片
                temp_dir = tempfile.mkdtemp()
                split_paths = []
                
                for i in range(split_count):
                    # 计算分割区域
                    top = int(i * actual_page_height)
                    # 确保最后一张图片包含到原始图片的底部，避免因浮点数精度问题导致的像素丢失
                    if i == split_count - 1:
                        bottom = height
                    else:
                        bottom = int((i + 1) * actual_page_height)
                    
                    # 分割图片
                    split_img = img.crop((0, top, width, bottom))
                    
                    # 生成分割后图片的文件名
                    base_name = os.path.splitext(os.path.basename(image_path))[0]
                    ext = os.path.splitext(image_path)[1]
                    split_filename = f"{base_name}_part{i+1}{ext}"
                    split_path = os.path.join(temp_dir, split_filename)
                    
                    # 保存分割后的图片
                    split_img.save(split_path)
                    split_paths.append(split_path)
                    print(f"    已保存分割部分: {split_filename}")
                
                return split_paths
            else:
                print(f"  高度差异比率 > {tolerANCE:.0%}，不满足分割条件，视为单张图片。")
                return [image_path]
                
    except Exception as e:
        print(f"处理图片时出错 {image_path}: {e}")
        # 出错时也返回原图，避免中断流程
        return [image_path]

def get_image_sort_key(filename):
    """
    提取文件名中的数字用于排序
    支持 page1.png, image_part1.png 等格式
    """
    # 提取所有数字
    numbers = re.findall(r'\d+', filename)
    if numbers:
        # 如果有多个数字，使用最后一个作为主要排序键，倒数第二个作为次要排序键
        if len(numbers) >= 2:
            return (int(numbers[-2]), int(numbers[-1]))
        else:
            return (int(numbers[0]), 0)
    return (0, 0)

# 获取当前脚本所在目录的路径
image_folder = os.getcwd()  # 获取当前工作目录
output_ppt = os.path.basename(image_folder) + ".pptx"  # 使用文件夹名称作为PPT的文件名

# 创建一个PPT文件
prs = Presentation()

# 获取所有图片文件 (注意：'JPG' 和 'JPEG' 都应包含)
image_files = [f for f in os.listdir(image_folder) if f.lower().endswith(('jpg', 'jpeg', 'png'))]

# 按数字进行排序
image_files.sort(key=get_image_sort_key)

# 获取幻灯片的宽度和高度 (单位: EMU - English Metric Units)
slide_width = prs.slide_width
slide_height = prs.slide_height

# 存储所有需要清理的临时文件路径
temp_paths_to_clean = []

print("开始处理图片...")

# 遍历图片文件
for image_file in image_files:
    img_path = os.path.join(image_folder, image_file)
    print(f"\n{'='*50}")
    print(f"正在处理: {image_file}")
    
    # 检测并分割图片（如果需要）
    split_paths = detect_and_split_vertical_images(img_path)
    
    print(f"分割结果: 获得 {len(split_paths)} 张图片")
    
    # 如果分割后的路径不是原路径，说明进行了分割，需要记录临时文件
    if len(split_paths) > 1 or (split_paths and split_paths[0] != img_path):
        # 只添加新生成的路径（即不是原始路径的）
        temp_paths_to_clean.extend([p for p in split_paths if p != img_path and p not in temp_paths_to_clean]) 
    
    # 将分割后的每张图片添加到PPT
    for idx, split_path in enumerate(split_paths):
        print(f"  添加到PPT: {os.path.basename(split_path)}")
        
        # 添加一个新的幻灯片 (使用标题和内容布局，但只添加图片)
        slide_layout = prs.slide_layouts[5]  # 选择一个空白的幻灯片布局 (标题和内容)
        slide = prs.slides.add_slide(slide_layout)

        # 设置图片的尺寸和位置（让图片铺满整个幻灯片）
        left = 0  # 图片左边距
        top = 0   # 图片上边距
        # 使用 slide_width 和 slide_height 确保图片填充整个幻灯片
        pic = slide.shapes.add_picture(split_path, left, top, width=slide_width, height=slide_height)

# 保存PPT文件
prs.save(output_ppt)

print(f"\nPPT已生成，保存路径：{os.path.abspath(output_ppt)}")

# 清理临时文件
print("\n清理临时文件...")
if temp_paths_to_clean:
    for temp_path in temp_paths_to_clean:
        try:
            os.remove(temp_path)
            print(f"  已删除临时文件: {temp_path}")
        except Exception as e:
             # 如果删除失败（例如文件不存在），打印警告但不中断
             print(f"  警告：清理临时文件失败 {temp_path}: {e}")

    # 尝试删除临时目录（如果为空）
    # 假设所有临时文件都在同一个目录下（由 tempfile.mkdtemp 创建）
    temp_dirs = set(os.path.dirname(p) for p in temp_paths_to_clean)
    for temp_dir in temp_dirs:
        try:
            os.rmdir(temp_dir)
            print(f"  已删除临时目录: {temp_dir}")
        except OSError as e:
            # 目录不为空或其他原因，忽略
            print(f"  信息：无法删除临时目录 {temp_dir} (可能非空): {e}")
        except Exception as e:
             print(f"  错误：尝试删除临时目录时发生意外错误 {temp_dir}: {e}")
else:
    print("  没有需要清理的临时文件。")

print("处理完成！")
